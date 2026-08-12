# -*- coding: utf-8 -*-
"""
Extrae el contenido de las presentaciones tipo A (VOCABULARIO, 11 slides) del
proyecto "Claudia Code" y lo vuelca a JSON + imágenes, para alimentar el HTML
de repaso.

Mapa slide -> bloque de la ficha de repaso:
    slide 1  -> palabra (word-box) + imagen "Dibujo" (portada)
    slide 2  -> articulemas (secuencia de fonemas)
    slide 3  -> tabla de sílabas
    slide 5  -> definición
    slide 9  -> 3 sinónimos + 3 antónimos, cada uno con pictograma
    slide 11 -> frase modelo + su imagen

Uso:  python extraer_pptx.py
"""
import hashlib
import io
import json
import re
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent))
from fonemas import fonemas as derivar_fonemas

PROYECTO = Path(r"C:\Users\santi\Documents\Claudia Code")
BASE = Path(r"C:\Users\santi\Documents\Repaso Programa Vocabulario")
PROGRAMA = PROYECTO / "📚 Programa Vocabulario y Lectura"
ARTICULEMAS = PROYECTO / "ARTICULEMAS"

# Carpeta del curso -> (código de curso, etapa, nombre legible)
CURSOS = {
    "INF 3 UNIDADES VOCABULARIO Y LECTURA": ("INF3", "infantil", "Infantil 3 años"),
    "INF 4 UNIDADES VOCABULARIO Y LECTURA": ("INF4", "infantil", "Infantil 4 años"),
    "INF 5 UNIDADES VOCABULARIO Y LECTURA": ("INF5", "infantil", "Infantil 5 años"),
    "1º EP UNIDADES VOCABULARIO Y LECTURA": ("1EP", "primaria", "1º Primaria"),
    "5º EP UNIDADES VOCABULARIO Y LECTURA": ("5EP", "primaria", "5º Primaria"),
}

# Frontera x (EMU) entre la columna de sinónimos y la de antónimos en slide 9
X_MEDIO = 4_660_825


def hash_bytes(b):
    return hashlib.sha1(b).hexdigest()


def mapa_articulemas():
    """hash del PNG -> letra/fonema, a partir de la carpeta ARTICULEMAS/."""
    m = {}
    for f in ARTICULEMAS.glob("_*_.png"):
        m[hash_bytes(f.read_bytes())] = f.stem.strip("_")
    return m


def texto(shape):
    if not shape.has_text_frame:
        return ""
    partes = []
    for p in shape.text_frame.paragraphs:
        t = "".join(r.text for r in p.runs)
        if t.strip():
            partes.append(t.strip())
    return "\n".join(partes)


def limpiar(s):
    s = s.replace("\u00a0", " ")
    s = re.sub(r"\s+", " ", s).strip()
    return s


def pictures(slide):
    return [sh for sh in slide.shapes if sh.shape_type == 13]


def guardar_img(shape, destino):
    """Guarda la imagen de un shape, aplanada a RGB y reescalada a un máximo."""
    blob = shape.image.blob
    im = Image.open(io.BytesIO(blob))
    if im.mode in ("RGBA", "LA", "P"):
        im = im.convert("RGBA")
        fondo = Image.new("RGB", im.size, (255, 255, 255))
        fondo.paste(im, (0, 0), im)
        im = fondo
    else:
        im = im.convert("RGB")
    im.thumbnail((520, 520), Image.LANCZOS)
    destino.parent.mkdir(parents=True, exist_ok=True)
    im.save(destino, "WEBP", quality=82, method=6)
    return destino.name


def extraer_silabas(slide):
    for sh in slide.shapes:
        if not sh.has_table:
            continue
        filas = sh.table.rows
        celdas = []
        for fila in filas:
            for c in fila.cells:
                celdas.append(limpiar(c.text))
        # la fila de palmadas trae emojis; la de sílabas, texto
        silabas = [c for c in celdas if c and not any(ord(ch) > 0x2000 for ch in c)]
        return silabas
    return []


def extraer_tipoA(ruta, curso, unidad, art_map, salida_img):
    prs = Presentation(str(ruta))
    slides = list(prs.slides)
    if len(slides) < 11:
        return None, f"solo {len(slides)} slides"

    ficha = {"curso": curso, "unidad": unidad, "archivo": ruta.name}

    # --- slide 1: palabra + dibujo -----------------------------------------
    s1 = slides[0]
    cajas = [sh for sh in s1.shapes if sh.has_text_frame and limpiar(texto(sh))]
    # el word-box es la caja de abajo (top mayor)
    wb = max(cajas, key=lambda sh: sh.top or 0)
    palabra = limpiar(texto(wb))
    ficha["palabra"] = palabra
    slug = re.sub(r"[^A-Za-z0-9]+", "", palabra.upper().translate(
        str.maketrans("ÁÉÍÓÚÜÑ", "AEIOUUN")))
    prefijo = f"{curso}_{unidad}_{slug}"

    pics = pictures(s1)
    ficha["img_dibujo"] = guardar_img(pics[0], salida_img / f"{prefijo}_dibujo.webp") if pics else None

    # --- slide 2: articulemas ----------------------------------------------
    # Las imágenes del pptx se reescribieron al construir los powers, así que su
    # hash ya no coincide con el set ARTICULEMAS/. La secuencia se deriva por
    # reglas grafema->fonema y se VALIDA contra el nº real de imágenes del pptx,
    # que sí es contenido ya revisado por la usuaria.
    seq = derivar_fonemas(palabra)
    n_pptx = len(pictures(slides[1]))
    ficha["fonemas"] = seq
    ficha["fonemas_pptx"] = n_pptx
    ficha["fonemas_ok"] = (len(seq) == n_pptx)

    # --- slide 3: sílabas ---------------------------------------------------
    ficha["silabas"] = extraer_silabas(slides[2])

    # --- slide 5: definición ------------------------------------------------
    s5 = slides[4]
    cand = [sh for sh in s5.shapes
            if sh.has_text_frame and limpiar(texto(sh))
            and (sh.width or 0) > 4_000_000 and (sh.top or 0) > 400_000]
    defi = ""
    for sh in sorted(cand, key=lambda s: s.top or 0):
        t = limpiar(texto(sh).replace("\n", " "))
        if t.upper() == palabra.upper():
            continue
        defi = t
        break
    ficha["definicion"] = defi

    # --- slide 9: sinónimos y antónimos ------------------------------------
    s9 = slides[8]
    etiquetas = [sh for sh in s9.shapes
                 if sh.has_text_frame and limpiar(texto(sh))
                 and 300_000 < (sh.width or 0) < 1_600_000
                 and (sh.top or 0) > 1_800_000]
    etiquetas.sort(key=lambda s: s.left or 0)
    sin_txt = [limpiar(texto(s)) for s in etiquetas if (s.left or 0) < X_MEDIO]
    ant_txt = [limpiar(texto(s)) for s in etiquetas if (s.left or 0) >= X_MEDIO]

    pics9 = sorted(pictures(s9), key=lambda s: s.left or 0)
    sin_pics = [p for p in pics9 if (p.left or 0) < X_MEDIO]
    ant_pics = [p for p in pics9 if (p.left or 0) >= X_MEDIO]

    def pareja(txts, pics, tipo):
        out = []
        for i, t in enumerate(txts):
            img = None
            if i < len(pics):
                img = guardar_img(pics[i], salida_img / f"{prefijo}_{tipo}{i+1}.webp")
            out.append({"texto": t, "img": img})
        return out

    ficha["sinonimos"] = pareja(sin_txt, sin_pics, "sin")
    ficha["antonimos"] = pareja(ant_txt, ant_pics, "ant")

    # --- slide 11: frase modelo --------------------------------------------
    s11 = slides[10]
    cajas11 = [sh for sh in s11.shapes
               if sh.has_text_frame and limpiar(texto(sh)) and (sh.top or 0) > 700_000]
    frase = ""
    for sh in cajas11:
        t = limpiar(texto(sh).replace("\n", " "))
        if t.upper() == palabra.upper() or "NIÑOS Y NIÑAS" in t.upper():
            continue
        frase = t
        break
    ficha["frase"] = frase
    pics11 = pictures(s11)
    ficha["img_frase"] = guardar_img(pics11[0], salida_img / f"{prefijo}_frase.webp") if pics11 else None

    return ficha, None


def main():
    art_map = mapa_articulemas()
    print(f"Articulemas conocidos: {len(art_map)}")

    resultados = {}
    incidencias = []

    for carpeta, (codigo, etapa, nombre) in CURSOS.items():
        raiz = PROGRAMA / carpeta
        if not raiz.is_dir():
            continue
        fichas = []
        for uni in sorted(raiz.iterdir(), key=lambda p: (len(p.name), p.name)):
            if not uni.is_dir():
                continue
            m = re.match(r"U(\d+)", uni.name)
            unidad = f"U{m.group(1)}" if m else uni.name
            for f in sorted(uni.glob("*.pptx")):
                # ojo: filtrar por "DIA" a secas se come palabras como ME-DIA-NO
                if "RECUPERACION" in f.name.upper() or re.search(r"_DIA[1-4]_", f.name.upper()):
                    continue
                salida_img = BASE / "assets" / "img" / codigo
                try:
                    ficha, err = extraer_tipoA(f, codigo, unidad, art_map, salida_img)
                except Exception as e:
                    incidencias.append(f"{f.name}: ERROR {e}")
                    continue
                if err:
                    incidencias.append(f"{f.name}: {err}")
                    continue
                ficha["etapa"] = etapa
                ficha["curso_nombre"] = nombre
                fichas.append(ficha)
                if not ficha["fonemas_ok"]:
                    incidencias.append(
                        f"{f.name}: fonemas regla={len(ficha['fonemas'])} "
                        f"vs pptx={ficha['fonemas_pptx']}")
                print(f"  {codigo} {unidad} {ficha['palabra']:<14} "
                      f"fon={len(ficha['fonemas'])}{'' if ficha['fonemas_ok'] else ' MISMATCH'} "
                      f"sin={len(ficha['sinonimos'])} ant={len(ficha['antonimos'])} "
                      f"def={'ok' if ficha['definicion'] else 'FALTA'} "
                      f"frase={'ok' if ficha['frase'] else 'FALTA'}")
        resultados[codigo] = fichas

    (BASE / "datos").mkdir(exist_ok=True)
    for codigo, fichas in resultados.items():
        with open(BASE / "datos" / f"{codigo}.json", "w", encoding="utf-8") as fh:
            json.dump(fichas, fh, ensure_ascii=False, indent=1)
        print(f"{codigo}: {len(fichas)} fichas")

    if incidencias:
        print("\nINCIDENCIAS:")
        for i in incidencias:
            print(" -", i)


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()
